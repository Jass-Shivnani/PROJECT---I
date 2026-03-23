"""
Extend parallax scrolling background from slide 23 to slide 38.

Strategy:
1. Apply gentle cumulative drift to blurred background images (Picture 34, 36)
   so they're not frozen at the same position across all slides.
2. Add 5 decorative circle images (extracted from slide 22's existing circles)
   to slides 24-38 with section-based positioning and small per-slide drift.
3. Use consistent shape names ("ParaDeco1"-"ParaDeco5") for Morph animation matching.
4. Place new shapes behind existing content in z-order.
5. Create an "entry slide" on slide 24 with circles off-screen so Morph
   smoothly animates them onto slide 25.

Parallax layers:
  SLOW   — Large blurred backgrounds drifting ~100-150K EMU/slide
  MEDIUM — Mid-size circles drifting ~200-400K EMU/slide
  FAST   — Small circles repositioning dramatically at section boundaries
"""

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import io
import math
import copy

pptx_path = r"docs\presentation\Capstone Review I.pptx"
prs = Presentation(pptx_path)

SLIDE_W = prs.slide_width   # 12192000
SLIDE_H = prs.slide_height  # 6858000

print(f"Slide dims: {SLIDE_W} x {SLIDE_H}")
print(f"Total slides: {len(prs.slides)}")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 1: Extract decorative circle image blobs from slide 22
# ═══════════════════════════════════════════════════════════════════════════════
slide_22 = prs.slides[21]  # 0-indexed
source_circles = {}
for shape in slide_22.shapes:
    if hasattr(shape, 'image') and shape.name in [
        'Picture 12', 'Picture 29', 'Picture 33',
        'Picture 39', 'Picture 41', 'Picture 43', 'Picture 45'
    ]:
        source_circles[shape.name] = {
            'blob': shape.image.blob,
            'width': shape.width,
            'height': shape.height,
        }
        print(f"  Extracted {shape.name}: {shape.width}x{shape.height}")

print(f"\nExtracted {len(source_circles)} circle images from slide 22")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 2: Apply gentle drift to blurred background images (slides 23-38)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n--- Applying background drift ---")

# Base positions (frozen values from slide 22)
bg_base = {
    'Picture 36': {'left': -4131617, 'top': -3941550},  # 9578170 x 8006808 EMU
    'Picture 34': {'left':  9305990, 'top': -2238550},  # 5734317 x 5748694 EMU
}

# Drift per slide — very gentle, creates subtle background movement
bg_drift_cfg = {
    'Picture 36': {'dl': -100000, 'dt':  40000},  # Drifts left and slightly down
    'Picture 34': {'dl': -150000, 'dt':  60000},  # Drifts left and slightly down
}

for slide_idx in range(22, 38):  # Slides 23-38 (0-indexed: 22-37)
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    offset = slide_num - 22  # 1 for slide 23, 16 for slide 38

    for shape in slide.shapes:
        if shape.name in bg_drift_cfg:
            cfg = bg_drift_cfg[shape.name]
            base = bg_base[shape.name]
            # Cumulative drift + sinusoidal wobble for organic feel
            new_left = int(
                base['left']
                + cfg['dl'] * offset
                + 70000 * math.sin(offset * 0.45)
            )
            new_top = int(
                base['top']
                + cfg['dt'] * offset
                + 50000 * math.cos(offset * 0.35)
            )
            old_left, old_top = shape.left, shape.top
            shape.left = new_left
            shape.top = new_top
            if offset <= 2 or offset >= 15:
                print(f"  Slide {slide_num} {shape.name}: "
                      f"({old_left},{old_top}) → ({new_left},{new_top})")

print("  ... (intermediate slides updated)")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 3: Define decorative circle layout across slides 24-38
# ═══════════════════════════════════════════════════════════════════════════════

# Image sources for each deco layer
deco_config = [
    {
        'name': 'ParaDeco1',
        'source': 'Picture 12',   # 1325563 x 1325563 (large circle)
        'speed': 'medium',
    },
    {
        'name': 'ParaDeco2',
        'source': 'Picture 29',   # 1185563 x 1185563
        'speed': 'medium',
    },
    {
        'name': 'ParaDeco3',
        'source': 'Picture 33',   # 862307 x 862307 (small)
        'speed': 'fast',
    },
    {
        'name': 'ParaDeco4',
        'source': 'Picture 39',   # 910008 x 910008 (small)
        'speed': 'fast',
    },
    {
        'name': 'ParaDeco5',
        'source': 'Picture 43',   # 1074874 x 1074874 (medium)
        'speed': 'medium',
    },
]

# Section definitions with slide ranges
# Each section has anchor positions for each deco circle
# Format: (base_left, base_top, drift_left/slide, drift_top/slide)
sections = {
    'entry':  {'slides': [24],       'label': 'Entry (off-screen for Morph)'},
    'A':      {'slides': [25,26,27,28], 'label': 'Block/Modular Diagrams'},
    'B':      {'slides': [29,30,31,32,33], 'label': 'Work Plan'},
    'C':      {'slides': [34,35,36,37,38], 'label': 'Conclusion/Refs/Thanks'},
}

# Position anchors for each deco in each section
# The 'entry' section places everything off-screen for slide 24
# (base_left, base_top, drift_left_per_slide, drift_top_per_slide)
positions = {
    'ParaDeco1': {
        'entry': (14500000, 1000000, 0, 0),      # off-screen right
        'A': (1500000, 900000, 120000, 40000),    # top-left area, drifting right/down
        'B': (9200000, 4500000, -80000, -50000),  # bottom-right, drifting left/up
        'C': (2200000, 5200000, 90000, -30000),   # bottom-left, drifting right
    },
    'ParaDeco2': {
        'entry': (15000000, 5500000, 0, 0),       # off-screen right
        'A': (10000000, 5000000, -100000, -60000), # bottom-right, drifting left
        'B': (1800000, 1000000, 80000, 70000),    # top-left, drifting right/down
        'C': (8200000, 900000, -70000, 50000),    # top-right, drifting left/down
    },
    'ParaDeco3': {
        'entry': (13500000, 3500000, 0, 0),       # off-screen right
        'A': (5500000, 5800000, -50000, -40000),  # bottom-center, drifting left/up
        'B': (10200000, 1800000, -120000, 60000), # top-right, drifting left/down
        'C': (1200000, 2800000, 70000, -50000),   # left-center, drifting right/up
    },
    'ParaDeco4': {
        'entry': (14000000, 6200000, 0, 0),       # off-screen right-bottom
        'A': (10500000, 2000000, -80000, 50000),  # right area, drifting left/down
        'B': (1000000, 5200000, 60000, -70000),   # bottom-left, drifting right/up
        'C': (6800000, 5500000, -60000, -40000),  # bottom-center, drifting left
    },
    'ParaDeco5': {
        'entry': (13000000, 800000, 0, 0),        # off-screen right-top
        'A': (3200000, 4200000, 70000, -50000),   # center-left, drifting right/up
        'B': (7200000, 700000, -90000, 80000),    # top-center, drifting left/down
        'C': (4200000, 3800000, 50000, -40000),   # center, gentle drift
    },
}

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 4: Add decorative circles to slides 24-38
# ═══════════════════════════════════════════════════════════════════════════════
print("\n--- Adding decorative circles ---")

def get_section(slide_num):
    """Get section key for a slide number."""
    if slide_num == 24:
        return 'entry'
    elif slide_num <= 28:
        return 'A'
    elif slide_num <= 33:
        return 'B'
    else:
        return 'C'

def get_section_start(section):
    """Get first slide number of a section."""
    return {'entry': 24, 'A': 25, 'B': 29, 'C': 34}[section]


added_count = 0

for slide_idx in range(23, 38):  # Slides 24-38 (0-indexed: 23-37)
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    section = get_section(slide_num)
    section_start = get_section_start(section)
    section_offset = slide_num - section_start  # 0, 1, 2, ...

    for deco in deco_config:
        deco_name = deco['name']
        src_name = deco['source']
        src = source_circles[src_name]

        pos = positions[deco_name][section]
        base_left, base_top, dl, dt = pos

        # Calculate position with drift
        left = int(base_left + dl * section_offset)
        top = int(base_top + dt * section_offset)

        # Add picture shape
        stream = io.BytesIO(src['blob'])
        pic = slide.shapes.add_picture(stream, left, top, src['width'], src['height'])
        pic.name = deco_name  # CRITICAL: Morph matches by name

        # Move shape behind existing content (insert at beginning of spTree)
        sp_tree = slide.shapes._spTree
        pic_element = pic._element
        sp_tree.remove(pic_element)
        # Insert after nvGrpSpPr and grpSpPr (first 2 children), before other shapes
        # Find the first sp/pic child to insert before
        insert_idx = 2  # After nvGrpSpPr and grpSpPr
        sp_tree.insert(insert_idx, pic_element)

        added_count += 1

    if slide_num in [24, 25, 29, 34, 38]:
        print(f"  Slide {slide_num} ({section}): Added 5 deco circles "
              f"(offset={section_offset})")

print(f"\nTotal decorative shapes added: {added_count}")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 5: Verify Morph transitions exist on all target slides
# ═══════════════════════════════════════════════════════════════════════════════
print("\n--- Verifying Morph transitions ---")

morph_count = 0
for slide_idx in range(22, 38):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    xml_str = etree.tostring(slide._element, pretty_print=False).decode()
    has_morph = 'morph' in xml_str.lower()
    if has_morph:
        morph_count += 1
    else:
        print(f"  ⚠ Slide {slide_num}: NO Morph transition!")

print(f"  {morph_count}/16 slides have Morph transitions ✓")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 6: Save the presentation
# ═══════════════════════════════════════════════════════════════════════════════
output_path = pptx_path.replace(".pptx", "_PARALLAX.pptx")
prs.save(output_path)
print(f"\n✓ Saved to {output_path}")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 7: Post-save verification — count images per slide
# ═══════════════════════════════════════════════════════════════════════════════
print("\n--- Post-save image count verification ---")
prs2 = Presentation(output_path)
for idx in range(21, 38):
    slide = prs2.slides[idx]
    slide_num = idx + 1
    pic_count = sum(1 for s in slide.shapes if hasattr(s, 'image'))
    # Count visible images
    vis_count = sum(
        1 for s in slide.shapes
        if hasattr(s, 'image')
        and s.left + s.width > 0 and s.left < SLIDE_W
        and s.top + s.height > 0 and s.top < SLIDE_H
    )
    bar = "█" * pic_count
    print(f"  Slide {slide_num:2d}: {pic_count:2d} total ({vis_count:2d} visible) {bar}")

# Check ParaDeco name consistency across consecutive slides
print("\n--- Morph name continuity check ---")
for slide_idx in range(23, 37):  # slides 24-37 → compare with next slide
    slide_a = prs2.slides[slide_idx]
    slide_b = prs2.slides[slide_idx + 1]
    names_a = {s.name for s in slide_a.shapes if s.name.startswith('ParaDeco')}
    names_b = {s.name for s in slide_b.shapes if s.name.startswith('ParaDeco')}
    common = names_a & names_b
    a_num = slide_idx + 1
    b_num = slide_idx + 2
    if len(common) < 5:
        print(f"  Slides {a_num}→{b_num}: {len(common)}/5 ParaDeco names match "
              f"(missing: {names_a - names_b})")
    else:
        print(f"  Slides {a_num}→{b_num}: 5/5 ParaDeco names match ✓")

print("\n✓ Parallax extension complete!")
