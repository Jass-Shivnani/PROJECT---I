"""Deep analysis of all slides 22-38 to plan parallax extension."""
from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r"docs\presentation\Capstone Review I.pptx"
prs = Presentation(pptx_path)

SLIDE_W = prs.slide_width   # 12192000 EMU
SLIDE_H = prs.slide_height  # 6858000 EMU

print(f"Slide dimensions: {SLIDE_W} x {SLIDE_H} EMU")
print(f"Total slides: {len(prs.slides)}\n")

# Collect all picture shapes across ALL slides
all_images = {}  # name -> {slide_idx: (left, top, width, height)}
slide_images = {}  # slide_idx -> list of (name, left, top, w, h)

for idx, slide in enumerate(prs.slides):
    slide_num = idx + 1
    slide_images[slide_num] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, 'image'):
                name = shape.name
                if name not in all_images:
                    all_images[name] = {}
                all_images[name][slide_num] = (shape.left, shape.top, shape.width, shape.height)
                slide_images[slide_num].append((name, shape.left, shape.top, shape.width, shape.height))

# Print slide-by-slide inventory for slides 22-38
print("=" * 80)
print("SLIDE-BY-SLIDE IMAGE INVENTORY (Slides 22-38)")
print("=" * 80)
for s in range(22, min(39, len(prs.slides) + 1)):
    imgs = slide_images.get(s, [])
    print(f"\n--- Slide {s} ({len(imgs)} images) ---")
    for name, l, t, w, h in sorted(imgs, key=lambda x: x[0]):
        # Check if on screen
        on_screen = (l + w > 0 and l < SLIDE_W and t + h > 0 and t < SLIDE_H)
        vis = "VISIBLE" if on_screen else "off-screen"
        print(f"  {name:20s} left={l:>10d} top={t:>10d} w={w:>9d} h={h:>9d}  [{vis}]")

# Identify the SLOW (blurred) images
print("\n" + "=" * 80)
print("SLOW/BLURRED BACKGROUND IMAGES (present on many slides, small deltas)")
print("=" * 80)

# From the analysis: Picture 31, Picture 34, Picture 36 are SLOW
slow_names = []
for name, positions in all_images.items():
    slides_on = sorted(positions.keys())
    if len(slides_on) < 4:
        continue
    # Calculate average delta per slide
    total_dl, total_dt, count = 0, 0, 0
    for i in range(len(slides_on) - 1):
        s1, s2 = slides_on[i], slides_on[i+1]
        gap = s2 - s1
        if gap == 0:
            continue
        dl = (positions[s2][0] - positions[s1][0]) / gap
        dt = (positions[s2][1] - positions[s1][1]) / gap
        total_dl += abs(dl)
        total_dt += abs(dt)
        count += 1
    if count > 0:
        avg_speed = (total_dl + total_dt) / count
        if avg_speed < 800000:  # SLOW threshold
            slow_names.append(name)
            print(f"  {name}: avg_speed={avg_speed:.0f}, appears on slides {slides_on}")
            # Show size
            sample = list(positions.values())[0]
            print(f"    Size: {sample[2]}x{sample[3]} EMU ({sample[2]/914400:.1f}x{sample[3]/914400:.1f} in)")

# Now identify which slides 23-38 are MISSING the slow bg images
print("\n" + "=" * 80)
print("MISSING BLURRED BG IMAGES ON SLIDES 23-38")
print("=" * 80)
for name in slow_names:
    positions = all_images[name]
    present = set(positions.keys())
    missing = []
    for s in range(23, min(39, len(prs.slides) + 1)):
        if s not in present:
            missing.append(s)
    if missing:
        print(f"  {name}: MISSING on slides {missing}")
        # Last known position
        last_slide = max(s for s in positions if s <= 22) if any(s <= 22 for s in positions) else max(positions)
        print(f"    Last position (slide {last_slide}): left={positions[last_slide][0]}, top={positions[last_slide][1]}")
    else:
        print(f"  {name}: present on all slides 23-38")

# Check which slides have Morph transitions
print("\n" + "=" * 80)
print("TRANSITION TYPES (slides 22-38)")
print("=" * 80)
from lxml import etree
for idx in range(21, min(38, len(prs.slides))):
    slide = prs.slides[idx]
    slide_num = idx + 1
    xml = etree.tostring(slide._element, pretty_print=False).decode()
    has_morph = 'morph' in xml.lower()
    has_transition = '<p:transition' in xml or '<mc:AlternateContent' in xml
    print(f"  Slide {slide_num}: morph={has_morph}, any_transition={has_transition}")

# Count total images per slide for slides 1-38
print("\n" + "=" * 80)
print("IMAGE COUNT PER SLIDE (all slides)")
print("=" * 80)
for s in range(1, len(prs.slides) + 1):
    count = len(slide_images.get(s, []))
    bar = "#" * count
    print(f"  Slide {s:2d}: {count:2d} images {bar}")
