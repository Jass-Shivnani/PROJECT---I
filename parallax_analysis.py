"""
Deep parallax analysis: track all background images across slides 1-22.
Map which images are blurred (slow) vs sharp (fast), and their position deltas.
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

prs = Presentation(r'docs\presentation\Capstone Review I.pptx')
SW = prs.slide_width
SH = prs.slide_height

# Track image positions across slides
# image_name -> {slide_num: (left, top, width, height)}
image_tracker = {}

for si in range(22):  # slides 1-22
    slide = prs.slides[si]
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
            name = shape.name
            if name not in image_tracker:
                image_tracker[name] = {}
            image_tracker[name][si + 1] = {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

# Find images that appear on multiple slides (these are the parallax ones)
parallax_images = {}
for name, appearances in image_tracker.items():
    if len(appearances) >= 2:
        parallax_images[name] = appearances

print(f"Total unique images: {len(image_tracker)}")
print(f"Parallax images (appear 2+ slides): {len(parallax_images)}")
print(f"Slide dims: {Emu(SW).inches:.2f}in x {Emu(SH).inches:.2f}in")
print(f"Slide dims EMU: {SW} x {SH}")
print()

# Analyze movement patterns
print("=== PARALLAX IMAGE MOVEMENT ===\n")
for name, appearances in sorted(parallax_images.items()):
    slides = sorted(appearances.keys())
    if len(slides) < 3:
        continue  # Need at least 3 appearances to see pattern
    
    print(f"--- {name} (appears on {len(slides)} slides: {slides[0]}-{slides[-1]}) ---")
    
    # Calculate per-slide deltas
    deltas_left = []
    deltas_top = []
    for i in range(1, len(slides)):
        prev_slide = slides[i-1]
        curr_slide = slides[i]
        slide_gap = curr_slide - prev_slide
        
        prev = appearances[prev_slide]
        curr = appearances[curr_slide]
        
        dl = (curr["left"] - prev["left"]) / slide_gap
        dt = (curr["top"] - prev["top"]) / slide_gap
        deltas_left.append(dl)
        deltas_top.append(dt)
        
        # Show position
        print(f"  Slide {prev_slide}→{curr_slide} (gap {slide_gap}): "
              f"left {prev['left']:>10} → {curr['left']:>10} (δ/slide={dl:>10.0f}), "
              f"top {prev['top']:>10} → {curr['top']:>10} (δ/slide={dt:>10.0f})")
    
    if deltas_left:
        avg_dl = sum(deltas_left) / len(deltas_left)
        avg_dt = sum(deltas_top) / len(deltas_top)
        # Classify speed
        speed = abs(avg_dl) + abs(avg_dt)
        category = "SLOW (blurred bg)" if speed < 500000 else "MEDIUM" if speed < 1500000 else "FAST (content)"
        print(f"  AVG δ/slide: left={avg_dl:.0f}, top={avg_dt:.0f} | Speed={speed:.0f} → {category}")
    print()


# Also check slides 22-24 to see where the pattern drops off
print("\n=== SLIDES 22-24 IMAGE INVENTORY ===")
for si in [21, 22, 23]:
    if si >= len(prs.slides):
        break
    slide = prs.slides[si]
    print(f"\nSlide {si+1}:")
    imgs = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
            imgs.append(f"  {shape.name}: left={shape.left}, top={shape.top}, "
                       f"w={shape.width}, h={shape.height}")
    for img in imgs:
        print(img)
    if not imgs:
        print("  (no images)")
